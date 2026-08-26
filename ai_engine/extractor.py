from __future__ import annotations

import json
import logging
import os
import re
import threading
import zipfile
from pathlib import Path

import settings

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    # Plain text / code
    ".txt", ".md", ".json", ".js", ".ts", ".tsx", ".jsx",
    ".py", ".html", ".css", ".csv", ".xml", ".yaml", ".yml",
    ".toml", ".rs", ".go", ".java", ".c", ".cpp", ".h",
    ".rb", ".sh", ".bat", ".env", ".sql", ".r",
    # Special parsers & Documents / Books
    ".pdf", ".docx", ".xlsx", ".pptx", ".ipynb", ".epub", ".rtf",
    # Images (OCR)
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff",
}

PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".json", ".js", ".ts", ".tsx", ".jsx",
    ".py", ".html", ".css", ".csv", ".xml", ".yaml", ".yml",
    ".toml", ".rs", ".go", ".java", ".c", ".cpp", ".h",
    ".rb", ".sh", ".bat", ".env", ".sql", ".r", ".rtf",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}

# Junk/lock files that should never waste embedding compute during content indexing
JUNK_FILENAMES = frozenset({
    "package-lock.json", "yarn.lock", "pnpm-lock.yaml", "cargo.lock",
    "poetry.lock", "gemfile.lock", "composer.lock", "tsconfig.tsbuildinfo",
    ".ds_store", "thumbs.db", "desktop.ini", "vocabulary.txt", "vocab.txt",
    "tokenizer.json", "tokenizer_config.json", "special_tokens_map.json",
})

JUNK_SUFFIXES = (
    ".min.js", ".min.css", ".bundle.js", ".map", ".d.ts",
    ".pyc", ".pyo", ".pyd", ".orig", ".rej", ".tmp", ".bak",
    ".swp", ".swo", ".log",
)

OCR_LANGS = [s for s in os.environ.get("LOCALMIND_OCR_LANGS", "en,tr").split(",") if s]
MAX_PDF_PAGES = int(os.environ.get("LOCALMIND_MAX_PDF_PAGES", "200"))


def ocr_enabled() -> bool:
    """OCR is opt-in: the models cost ~500MB of RAM and seconds per image."""
    return bool(settings.get("ocr"))


def is_junk_file(path: str | Path) -> bool:
    """Check if file is a build artifact, lockfile, minified bundle, or junk."""
    name = os.path.basename(path).lower()
    if name in JUNK_FILENAMES:
        return True
    if any(name.endswith(sfx) for sfx in JUNK_SUFFIXES):
        return True
    return False


def extract_text(file_path: str, max_size_mb: float = 50) -> str | None:
    """Extract text content from a file with C++ optimized parsers."""
    path = Path(file_path)
    name = path.name.lower()
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS or is_junk_file(name):
        return None
    if ext in IMAGE_EXTENSIONS and not ocr_enabled():
        return None

    try:
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > max_size_mb:
            logger.debug("Skipping %s (%.1f MB > %.1f MB limit)", file_path, size_mb, max_size_mb)
            return None

        if ext in PLAIN_TEXT_EXTENSIONS:
            return _read_plain_text(path)
        elif ext == ".pdf":
            return _read_pdf(path)
        elif ext == ".epub":
            return _read_epub(path)
        elif ext == ".docx":
            return _read_docx(path)
        elif ext == ".xlsx":
            return _read_xlsx(path)
        elif ext == ".pptx":
            return _read_pptx(path)
        elif ext == ".ipynb":
            return _read_ipynb(path)
        elif ext in IMAGE_EXTENSIONS:
            return _read_image_ocr(path)
    except Exception as e:
        logger.warning("Failed to extract text from %s: %s", file_path, e)
    return None


def _read_plain_text(path: Path) -> str | None:
    try:
        size = path.stat().st_size
        if size == 0:
            return None
        if size > 64 * 1024:
            # Memory-mapped read for larger files avoids redundant full-buffer copies
            import mmap
            with open(path, "rb") as f:
                with mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
                    for enc in ("utf-8", "latin-1", "cp1252", "utf-16"):
                        try:
                            return mm.read().decode(enc).strip() or None
                        except (UnicodeDecodeError, ValueError):
                            mm.seek(0)
                            continue
        else:
            for encoding in ("utf-8", "latin-1", "cp1252", "utf-16"):
                try:
                    text = path.read_text(encoding=encoding)
                    return text.strip() if text.strip() else None
                except (UnicodeDecodeError, ValueError):
                    continue
    except Exception:
        pass
    return None


def _read_epub(path: Path) -> str | None:
    """Extract readable text from an EPUB ebook using built-in zipfile."""
    try:
        texts = []
        with zipfile.ZipFile(path, "r") as z:
            for filename in z.namelist():
                if filename.lower().endswith((".xhtml", ".html", ".htm")):
                    raw = z.read(filename).decode("utf-8", errors="ignore")
                    clean = re.sub(r"<[^>]+>", " ", raw)
                    clean = " ".join(clean.split())
                    if clean:
                        texts.append(clean)
        return "\n\n".join(texts).strip() if texts else None
    except Exception as e:
        logger.debug("Failed to read epub %s: %s", path, e)
        return None


def _read_pdf(path: Path, max_pages: int = MAX_PDF_PAGES) -> str | None:
    """Extract PDF text using C++ PyMuPDF with instant C-heap glyph/pixmap cache shrinking."""
    # 1. Ultra-fast PyMuPDF (MuPDF C++ engine)
    doc = None
    try:
        import pymupdf
        doc = pymupdf.open(str(path))
        texts = []
        n_pages = min(len(doc), max_pages)
        for i in range(n_pages):
            page = doc[i]
            t = page.get_text()
            if t and t.strip():
                texts.append(t.strip())
        combined = "\n\n".join(texts).strip()
        if combined:
            return combined
    except Exception as e:
        logger.debug("PyMuPDF extraction fallback for %s: %s", path, e)
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass
        try:
            import pymupdf
            pymupdf.fitz.TOOLS.store_shrink(100)
        except Exception:
            pass

    # 2. Fallback to pypdfium2 C++ parser
    pdf = None
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(str(path))
        texts = []
        n_pages = min(len(pdf), max_pages)
        for i in range(n_pages):
            page = pdf[i]
            textpage = page.get_textpage()
            text = textpage.get_text_range()
            if text and text.strip():
                texts.append(text.strip())
        combined = "\n\n".join(texts).strip()
        if combined:
            return combined
    except Exception as e:
        logger.debug("pypdfium2 extraction fallback for %s: %s", path, e)
    finally:
        if pdf is not None:
            try:
                pdf.close()
            except Exception:
                pass

    # 3. Fallback to pdfplumber if needed
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf_plumb:
            for i, page in enumerate(pdf_plumb.pages):
                if i >= max_pages:
                    break
                try:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
                finally:
                    page.flush_cache()
                    try:
                        page.get_textmap.cache_clear()
                    except Exception:
                        pass
        combined = "\n".join(texts).strip()
        if combined:
            return combined
    except Exception:
        pass

    # 4. Fallback: OCR only when explicitly enabled (scanned PDFs)
    if ocr_enabled():
        return _read_image_ocr(path)
    return None


def _read_docx(path: Path) -> str | None:
    """Extract text from .docx via zero-overhead zipfile XML streaming (no heavy lxml DOM)."""
    try:
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(path, "r") as z:
            if "word/document.xml" in z.namelist():
                with z.open("word/document.xml") as xml_f:
                    tree = ET.parse(xml_f)
                    root = tree.getroot()
                    texts = [node.text for node in root.iter() if node.tag.endswith("}t") and node.text]
                    if texts:
                        return " ".join(texts).strip()
    except Exception:
        pass

    # Fallback to python-docx
    try:
        from docx import Document
        doc = Document(str(path))
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(texts).strip() if texts else None
    except Exception as e:
        logger.debug("docx extraction failed for %s: %s", path, e)
        return None


def _read_xlsx(path: Path) -> str | None:
    wb = None
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames[:5]:  # limit to first 5 sheets
            ws = wb[sheet]
            texts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(max_row=500, values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    texts.append(row_text)
        combined = "\n".join(texts).strip()
        return combined if combined else None
    except ImportError:
        logger.debug("openpyxl not installed, skipping .xlsx")
        return None
    except Exception as e:
        logger.debug("xlsx extraction error for %s: %s", path, e)
        return None
    finally:
        if wb is not None:
            try:
                wb.close()
            except Exception:
                pass


def _read_pptx(path: Path) -> str | None:
    """Extract text from .pptx via lightweight zipfile slide parsing."""
    try:
        import xml.etree.ElementTree as ET
        texts = []
        with zipfile.ZipFile(path, "r") as z:
            slide_files = sorted([n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")])
            for slide_num, slide_name in enumerate(slide_files, 1):
                with z.open(slide_name) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    slide_texts = [node.text for node in root.iter() if node.tag.endswith("}t") and node.text and node.text.strip()]
                    if slide_texts:
                        texts.append(f"[Slide {slide_num}]")
                        texts.extend(slide_texts)
        if texts:
            return "\n".join(texts).strip()
    except Exception:
        pass

    # Fallback to python-pptx
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                texts.append(f"[Slide {slide_num}]")
                texts.extend(slide_texts)
        combined = "\n".join(texts).strip()
        return combined if combined else None
    except ImportError:
        logger.debug("python-pptx not installed, skipping .pptx")
        return None


def _read_ipynb(path: Path) -> str | None:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        texts = []
        for cell in data.get("cells", []):
            source = "".join(cell.get("source", []))
            if source.strip():
                cell_type = cell.get("cell_type", "code")
                texts.append(f"[{cell_type}]")
                texts.append(source.strip())
        combined = "\n".join(texts).strip()
        return combined if combined else None
    except Exception:
        return None


_ocr_available: bool | None = None
_ocr_reader = None
_ocr_lock = threading.Lock()


def _get_ocr_reader():
    """Build the easyocr Reader once and reuse it."""
    global _ocr_available, _ocr_reader
    if _ocr_reader is not None:
        return _ocr_reader
    if _ocr_available is False:
        return None

    with _ocr_lock:
        if _ocr_reader is not None:
            return _ocr_reader
        if _ocr_available is False:
            return None
        try:
            import easyocr
            _ocr_reader = easyocr.Reader(OCR_LANGS, gpu=False, verbose=False)
            _ocr_available = True
            logger.info("OCR reader initialized (langs=%s)", OCR_LANGS)
        except ImportError:
            _ocr_available = False
            logger.debug("easyocr not installed, OCR disabled")
        except Exception as e:
            _ocr_available = False
            logger.warning("OCR reader init failed, OCR disabled: %s", e)
    return _ocr_reader


def release_ocr_reader() -> None:
    """Drop the OCR models from memory."""
    global _ocr_reader
    with _ocr_lock:
        if _ocr_reader is not None:
            _ocr_reader = None
            import gc
            gc.collect()
            logger.info("OCR reader released")


def _read_image_ocr(path: Path) -> str | None:
    if not ocr_enabled():
        return None

    reader = _get_ocr_reader()
    if reader is None:
        return None

    try:
        with _ocr_lock:
            results = reader.readtext(str(path), detail=1, paragraph=False)
        texts = [r[1] for r in results if r[1].strip()]
        combined = " ".join(texts).strip()
        return combined if combined else None
    except Exception as e:
        logger.warning("OCR failed for %s: %s", path, e)
        return None
